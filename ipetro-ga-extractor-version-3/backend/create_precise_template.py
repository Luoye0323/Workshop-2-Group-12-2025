import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_precise_template(output_file: str):
    """
    Creates a high-fidelity Inspection Plan template matching the user's requirements.
    Key features:
    - Exact table layout
    - Specific colors (Green Risk, Yellow Group, Grey Headers)
    - Placeholders for data filling
    - Blank sections where requested
    """
    # 1. Initialize Presentation (A4 Landscape)
    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)

    # Use blank layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Set slide background to white (remove blue default)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

    # --- Helper Functions ---
    def set_cell(
        cell,
        text,
        font_size=9,
        bold=False,
        align=PP_ALIGN.CENTER,
        bg_color=None,
        font_name="Arial",
        font_color=None,
    ):
        # 1. Clear existing content (leaves one empty paragraph)
        cell.text = ""

        # 2. Split text into lines to handle newlines explicitly
        lines = str(text).split("\n")

        for i, line in enumerate(lines):
            # Use the first paragraph if it's the first line, otherwise add new
            paragraph = (
                cell.text_frame.paragraphs[0]
                if i == 0
                else cell.text_frame.add_paragraph()
            )
            paragraph.alignment = align

            # Add run and apply style to the run directly
            run = paragraph.add_run()
            run.text = line
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = font_name

            if font_color:
                run.font.color.rgb = font_color

        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

        # Add thin black borders to all sides
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.util import Pt as BorderPt

        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        # Define border properties for all sides
        for border_name in ["lnL", "lnR", "lnT", "lnB"]:
            ln = OxmlElement("a:" + border_name)
            ln.set("w", "12700")  # Width in EMUs (12700 = 1pt)
            ln.set("cap", "flat")
            ln.set("cmpd", "sng")
            ln.set("algn", "ctr")

            solidFill = OxmlElement("a:solidFill")
            srgbClr = OxmlElement("a:srgbClr")
            srgbClr.set("val", "000000")  # Black color
            solidFill.append(srgbClr)
            ln.append(solidFill)

            prstDash = OxmlElement("a:prstDash")
            prstDash.set("val", "solid")
            ln.append(prstDash)

            tcPr.append(ln)

    def merge_cells(table, r1, c1, r2, c2, text=None, **kwargs):
        cell = table.cell(r1, c1)
        cell.merge(table.cell(r2, c2))
        if text:
            set_cell(cell, text, **kwargs)
        return cell

    # Colors
    GREY_HEADER = RGBColor(217, 217, 217)
    YELLOW_BG = RGBColor(255, 255, 153)
    GREEN_RISK = RGBColor(0, 176, 80)
    LIGHT_GREEN_BG = RGBColor(235, 241, 222)
    WHITE = RGBColor(255, 255, 255)

    # --- 2. HEADER SECTION ---
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(10.69), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "INSPECTION PLAN (ON-STREAM)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.name = "Emirates"
    p.alignment = PP_ALIGN.CENTER

    # General Info Table
    # Rows: 2, Cols: 5
    left, top, width, height = Inches(0.5), Inches(0.7), Inches(10.69), Inches(0.8)
    header_table = slide.shapes.add_table(2, 5, left, top, width, height).table

    # Columns: Desc(4.0), Tag(2.0), PMT(2.0), Rev(1.0), Logo(1.69)
    header_table.columns[0].width = Inches(4.0)
    header_table.columns[1].width = Inches(2.0)
    header_table.columns[2].width = Inches(2.0)
    header_table.columns[3].width = Inches(1.0)
    header_table.columns[4].width = Inches(1.69)

    # Row 0
    merge_cells(header_table, 0, 0, 0, 2, "GENERAL INFORMATION", bold=True)
    header_table.cell(0, 3).merge(header_table.cell(0, 4))
    set_cell(
        header_table.cell(0, 3), "PLANT UNIT: {{UNIT}}", bold=True, align=PP_ALIGN.LEFT
    )

    # Row 1
    set_cell(header_table.cell(1, 0), "Description : {{DESC}}", align=PP_ALIGN.LEFT)
    set_cell(header_table.cell(1, 1), "Tag. No : {{TAG}}", align=PP_ALIGN.LEFT)
    set_cell(header_table.cell(1, 2), "PMT. No : {{PMT}}", align=PP_ALIGN.LEFT)
    set_cell(header_table.cell(1, 3), "Rev 00", align=PP_ALIGN.CENTER)
    # Logo handling
    # Calculate position for logo (Row 1, Col 4)
    # Left = 0.5 + 4.0 + 2.0 + 2.0 + 1.0 = 9.5 inches
    logo_left = Inches(9.85)
    # Top = 0.7 (table top) + 0.4 (approx row 0 height)
    logo_top = Inches(1.14)
    # Dimensions
    logo_width = Inches(1.05)
    logo_height = Inches(0.32)

    # Try to find the logo
    # Use PNG with transparency preserved by sips conversion
    possible_logo_paths = [
        "public/ipetro-logo-transparent.png",  # Transparent PNG (best option)
        os.path.join(
            os.path.dirname(__file__), "public", "ipetro-logo-transparent.png"
        ),
        "public/ipetro-logo.svg.png",  # Fallback
        os.path.join(os.path.dirname(__file__), "public", "ipetro-logo.svg"),
    ]

    logo_inserted = False
    for logo_path in possible_logo_paths:
        if os.path.exists(logo_path):
            try:
                # Add picture centered in the cell area
                # For SVG, newer PPT versions handle it. If safe, python-pptx accepts it.
                # First insert with restricted width
                pic = slide.shapes.add_picture(
                    logo_path, logo_left, logo_top, width=logo_width
                )

                # Check if height exceeds the row height
                if pic.height > logo_height:
                    # Scale down to fit height
                    aspect_ratio = pic.width / pic.height
                    pic.height = int(logo_height)
                    pic.width = int(aspect_ratio * pic.height)

                logo_inserted = True
                break
            except Exception as e:
                print(f"Failed to insert logo {logo_path}: {e}")

    if not logo_inserted:
        # Fallback to text if no logo found
        set_cell(
            header_table.cell(1, 4),
            "iPETRO",
            align=PP_ALIGN.CENTER,
            bold=True,
            font_size=14,
        )

    # --- 3. MAIN EQUIPMENT TABLE ---
    # Rows: 5 (Header1, Header2, Data1, Data2, Data3)
    # Cols: 11
    top_y = Inches(1.6)
    rows, cols = 5, 11
    main_table = slide.shapes.add_table(
        rows, cols, Inches(0.5), top_y, Inches(10.69), Inches(2.0)
    ).table

    # Widths
    widths = [1.0, 1.5, 1.2, 0.6, 0.6, 0.6, 0.8, 0.6, 0.6, 1.0, 2.19]
    for i, w in enumerate(widths):
        main_table.columns[i].width = Inches(w)

    # Headers (Row 0 & 1)
    # Vertical Merges
    headers = [
        (0, "FLUID"),
        (1, "COMPONENT"),
        (2, "DESIGN CODE"),
        (6, "INSULATED\n(Y/N)"),
        (9, "CURRENT\nRISK\nRATING"),
    ]
    for c, text in headers:
        merge_cells(
            main_table,
            0,
            c,
            1,
            c,
            text,
            bold=True,
            font_color=RGBColor(0, 0, 0),
            bg_color=GREY_HEADER,
        )
    # Horizontal Merges
    merge_cells(
        main_table,
        0,
        3,
        0,
        5,
        "MATERIAL",
        bold=True,
        font_color=RGBColor(0, 0, 0),
        bg_color=GREY_HEADER,
    )
    merge_cells(
        main_table,
        0,
        7,
        0,
        8,
        "OP. PARAMETER",
        bold=True,
        font_color=RGBColor(0, 0, 0),
        bg_color=GREY_HEADER,
    )

    # Damage Mech Header (Yellow) - Merged across 2 rows
    dm_cell = merge_cells(main_table, 0, 10, 1, 10)
    # Custom text run for multi-format if needed, but simple text for now
    set_cell(
        dm_cell,
        "Corrosion Group :\nDamage Mechanism\nSusceptible:",
        bold=True,
        align=PP_ALIGN.LEFT,
        bg_color=YELLOW_BG,
        font_color=RGBColor(0, 0, 0),
    )

    # Subheaders (Row 1)
    sub_headers = {3: "TYPE", 4: "SPEC", 5: "GR", 7: "T\n(C)", 8: "P\n(Mpa)"}
    for c, text in sub_headers.items():
        set_cell(
            main_table.cell(1, c),
            text,
            bold=True,
            bg_color=GREY_HEADER,
            font_color=RGBColor(0, 0, 0),
        )

    # Data Rows (2, 3, 4) - Pre-filled KEY placeholders
    # Row 2: Top Head
    set_cell(main_table.cell(2, 0), "{{FLUID_TOP}}")
    set_cell(main_table.cell(2, 1), "Top Head")
    set_cell(main_table.cell(2, 2), "{{CODE_TOP}}")
    set_cell(main_table.cell(2, 3), "{{MAT_TOP_TYPE}}")
    set_cell(main_table.cell(2, 4), "{{MAT_TOP_SPEC}}")
    set_cell(main_table.cell(2, 5), "{{MAT_TOP_GRADE}}")
    set_cell(main_table.cell(2, 6), "{{INSUL}}")
    set_cell(main_table.cell(2, 7), "{{OP_TEMP}}")
    set_cell(main_table.cell(2, 8), "{{OP_PRESS}}")

    # Row 3: Shell
    set_cell(main_table.cell(3, 0), "{{FLUID_SHELL}}")
    set_cell(main_table.cell(3, 1), "Shell")
    set_cell(main_table.cell(3, 2), "{{CODE_SHELL}}")
    set_cell(main_table.cell(3, 3), "{{MAT_SHELL_TYPE}}")
    set_cell(main_table.cell(3, 4), "{{MAT_SHELL_SPEC}}")
    set_cell(main_table.cell(3, 5), "{{MAT_SHELL_GRADE}}")
    set_cell(main_table.cell(3, 6), "{{INSUL}}")
    set_cell(main_table.cell(3, 7), "{{OP_TEMP}}")
    set_cell(main_table.cell(3, 8), "{{OP_PRESS}}")

    # Row 4: Bottom Head
    set_cell(main_table.cell(4, 0), "{{FLUID_BOT}}")
    set_cell(main_table.cell(4, 1), "Bottom Head")
    set_cell(main_table.cell(4, 2), "{{CODE_BOT}}")
    set_cell(main_table.cell(4, 3), "{{MAT_BOT_TYPE}}")
    set_cell(main_table.cell(4, 4), "{{MAT_BOT_SPEC}}")
    set_cell(main_table.cell(4, 5), "{{MAT_BOT_GRADE}}")
    set_cell(main_table.cell(4, 6), "{{INSUL}}")
    set_cell(main_table.cell(4, 7), "{{OP_TEMP}}")
    set_cell(main_table.cell(4, 8), "{{OP_PRESS}}")

    # Merge RISK column (Rows 2-4)
    risk_cell = merge_cells(
        main_table,
        2,
        9,
        4,
        9,
        "LOW",
        bold=True,
        bg_color=GREEN_RISK,
        font_color=RGBColor(0, 0, 0),
    )

    # Merge DAM MECH column (Rows 2-4) - LEFT BLANK as requested
    dm_data_cell = merge_cells(
        main_table,
        2,
        10,
        4,
        10,
        "",
        align=PP_ALIGN.LEFT,
        bg_color=YELLOW_BG,
        font_color=RGBColor(0, 0, 0),
    )

    # --- 4. INSPECTION TABLE (Bottom Left) ---
    insp_top = Inches(4.65)
    insp_table = slide.shapes.add_table(
        3, 3, Inches(0.5), insp_top, Inches(5.0), Inches(2.5)
    ).table

    # Cols
    insp_table.columns[0].width = Inches(1.2)
    insp_table.columns[1].width = Inches(2.3)
    insp_table.columns[2].width = Inches(1.5)

    # Headers
    headers = ["INSPECTION\nMETHOD", "INSPECTION COVERAGE", "DAMAGE MECHANISM"]
    for i, h in enumerate(headers):
        set_cell(
            insp_table.cell(0, i),
            h,
            bold=True,
            bg_color=GREY_HEADER,
            font_color=RGBColor(0, 0, 0),
        )

    # Rows 1 & 2 - LEFT BLANK as requested (user said "content could be different")
    # We just create the cells with light green bg
    for r in [1, 2]:
        for c in [0, 1, 2]:
            set_cell(insp_table.cell(r, c), "", bg_color=LIGHT_GREEN_BG)

    # --- 5. DRAWING PLACEHOLDER (Bottom Right) ---
    left, top = Inches(5.8), Inches(3.8)
    width, height = Inches(5.39), Inches(3.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = "[EQUIPMENT DRAWING]"
    shape.fill.background()
    shape.line.color.rgb = RGBColor(180, 180, 180)

    # --- 6. FOOTER ---
    footer_top = Inches(7.5)
    footer_h = Inches(0.6)
    footer_table = slide.shapes.add_table(
        2, 3, Inches(0.5), footer_top, Inches(10.69), footer_h
    ).table

    footer_headers = [
        "Prepared by",
        "Review and Approved by",
        "Acknowledge and Acc. by",
    ]
    for i, h in enumerate(footer_headers):
        set_cell(
            footer_table.cell(0, i), h, bold=True, font_size=8, align=PP_ALIGN.LEFT
        )
        set_cell(
            footer_table.cell(1, i),
            "Date",
            bold=True,
            font_size=8,
            align=PP_ALIGN.CENTER,
        )

    # Save
    prs.save(output_file)
    print(f"Created template: {output_file}")
    return output_file


if __name__ == "__main__":
    # Save to the specific templates folder
    output_path = os.path.join("templates", "inspection_plan_template.pptx")
    # Ensure dir exists
    os.makedirs("templates", exist_ok=True)
    create_precise_template(output_path)
