import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_inspection_plan_template(output_file: str, data: dict = None):
    """
    Recreates the Inspection Plan template exactly using python-pptx tables.
    Fills with data if provided, otherwise uses placeholders like {{KEY}}.
    """
    if data is None:
        data = {}

    # 1. Initialize Presentation with A4 Landscape-ish dimensions
    prs = Presentation()
    prs.slide_width = Inches(11.69)  # A4 Landscape width
    prs.slide_height = Inches(8.27)  # A4 Landscape height
    
    # Use a blank layout
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)

    # --- Helper Functions for Styling ---
    def set_cell_text(cell, text, font_size=9, bold=False, align=PP_ALIGN.CENTER, bg_color=None):
        cell.text = str(text)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.alignment = align
        paragraph.font.name = 'Arial'
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

    def merge_cells(table, r1, c1, r2, c2, text=None):
        cell = table.cell(r1, c1)
        cell.merge(table.cell(r2, c2))
        if text:
            set_cell_text(cell, text, bold=True)
        return cell

    # --- 2. HEADER SECTION ---
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10.69), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "INSPECTION PLAN (ON-STREAM)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    # Top Info Table (Description, Tag, PMT, etc.)
    # Rows: 2, Cols: 5 (Description | Tag | PMT | Rev | Logo area)
    left, top, width, height = Inches(0.5), Inches(0.7), Inches(10.69), Inches(0.8)
    header_table = slide.shapes.add_table(2, 5, left, top, width, height).table
    
    # Set Column Widths for Header
    header_table.columns[0].width = Inches(4.0) # Description
    header_table.columns[1].width = Inches(2.0) # Tag
    header_table.columns[2].width = Inches(2.0) # PMT
    header_table.columns[3].width = Inches(1.0) # Rev
    header_table.columns[4].width = Inches(1.69)# Logo

    # Row 0: General Info Header
    merge_cells(header_table, 0, 0, 0, 2, "GENERAL INFORMATION")
    header_table.cell(0, 3).merge(header_table.cell(0, 4))
    set_cell_text(header_table.cell(0, 3), f"PLANT UNIT: {data.get('plant_unit', '{{UNIT}}')}", bold=True, align=PP_ALIGN.LEFT)

    # Row 1: Specific Data
    set_cell_text(header_table.cell(1, 0), f"Description : {data.get('description', '{{DESC}}')}", align=PP_ALIGN.LEFT)
    set_cell_text(header_table.cell(1, 1), f"Tag. No : {data.get('tag_no', '{{TAG}}')}", align=PP_ALIGN.LEFT)
    set_cell_text(header_table.cell(1, 2), f"PMT. No : {data.get('pmt_no', '{{PMT}}')}", align=PP_ALIGN.LEFT)
    set_cell_text(header_table.cell(1, 3), "Rev 00", align=PP_ALIGN.CENTER)
    set_cell_text(header_table.cell(1, 4), "[LOGO PLACEHOLDER]", align=PP_ALIGN.CENTER)

    # --- 3. MAIN EQUIPMENT TABLE (The Grid) ---
    # Rows: 5 (Header1, Header2, Data1, Data2, Data3)
    # Cols: 11
    top_y = Inches(1.6)
    rows, cols = 5, 11
    main_table = slide.shapes.add_table(rows, cols, Inches(0.5), top_y, Inches(10.69), Inches(2.0)).table

    # Column Widths (Approximate based on image)
    widths = [1.0, 1.5, 1.2, 0.6, 0.6, 0.6, 0.8, 0.6, 0.6, 1.0, 2.19]
    for i, w in enumerate(widths):
        main_table.columns[i].width = Inches(w)

    # -- Headers Row 0 & 1 Setup --
    grey_color = RGBColor(217, 217, 217)
    
    # Vertical Merges for Main Headers
    headers = [
        (0, 0, "FLUID"), (0, 1, "COMPONENT"), (0, 2, "DESIGN CODE"),
        (0, 6, "INSULATED\n(Y/N)"), (0, 9, "CURRENT\nRISK\nRATING")
    ]
    for c, text in [(0, "FLUID"), (1, "COMPONENT"), (2, "DESIGN CODE"), (6, "INSULATED\n(Y/N)"), (9, "CURRENT\nRISK\nRATING")]:
        cell = merge_cells(main_table, 0, c, 1, c, text)
        set_cell_text(cell, text, bold=True, bg_color=grey_color)

    # Horizontal Merges for Group Headers
    mat_header = merge_cells(main_table, 0, 3, 0, 5, "MATERIAL")
    set_cell_text(mat_header, "MATERIAL", bold=True, bg_color=grey_color)
    
    op_header = merge_cells(main_table, 0, 7, 0, 8, "OP. PARAMETER")
    set_cell_text(op_header, "OP. PARAMETER", bold=True, bg_color=grey_color)

    # Damage Mechanism Header (Yellow, large merge)
    dm_header = merge_cells(main_table, 0, 10, 1, 10, "")
    set_cell_text(dm_header, "Corrosion Group :\nDamage Mechanism\nSusceptible:", bold=True, align=PP_ALIGN.LEFT, bg_color=RGBColor(255, 255, 153))

    # Sub-headers (Row 1)
    sub_headers = {3: "TYPE", 4: "SPEC", 5: "GR", 7: "T\n(°C)", 8: "P\n(Mpa)"}
    for col_idx, text in sub_headers.items():
        set_cell_text(main_table.cell(1, col_idx), text, bold=True, bg_color=grey_color)

    # -- Data Filling (Rows 2, 3, 4) --
    # Example components list to fill rows
    components = data.get('components', [
        {'name': 'Top Head', 'code': 'ASME VIII DIV 1', 'mat': 'SS'},
        {'name': 'Shell', 'code': 'ASME VIII DIV 1', 'mat': 'SS'},
        {'name': 'Bottom Head', 'code': 'ASME VIII DIV 1', 'mat': 'SS'}
    ])

    # Fill components (loop through max 3 rows for this template)
    green_color = RGBColor(0, 176, 80)
    for i, comp in enumerate(components):
        r = i + 2
        if r >= 5: break
        set_cell_text(main_table.cell(r, 1), comp.get('name', '{{COMP}}'))
        set_cell_text(main_table.cell(r, 2), comp.get('code', '{{CODE}}'))
        set_cell_text(main_table.cell(r, 3), comp.get('mat', '{{MAT}}'))
        # Placeholder for other cols
        set_cell_text(main_table.cell(r, 0), "{{FLUID}}") 
        
    # Merge Risk Column for data rows (assuming single risk for vessel)
    risk_cell = merge_cells(main_table, 2, 9, 4, 9, "LOW")
    set_cell_text(risk_cell, "LOW", bold=True, bg_color=green_color)

    # Merge Damage Mech Column for data rows
    dm_data_cell = merge_cells(main_table, 2, 10, 4, 10, "")
    dm_text = "Internal\n• Cooling Water Corrosion\n\nExternal\n• General Corrosion"
    set_cell_text(dm_data_cell, data.get('damage_mechanisms', dm_text), align=PP_ALIGN.LEFT, bg_color=RGBColor(255, 255, 153))


    # --- 4. INSPECTION TABLE (Bottom Left) ---
    insp_top = Inches(4.0)
    insp_table = slide.shapes.add_table(3, 3, Inches(0.5), insp_top, Inches(5.0), Inches(2.5)).table
    
    # Headers
    headers = ["INSPECTION\nMETHOD", "INSPECTION COVERAGE", "DAMAGE MECHANISM"]
    insp_table.columns[0].width = Inches(1.2)
    insp_table.columns[1].width = Inches(2.3)
    insp_table.columns[2].width = Inches(1.5)

    for i, h in enumerate(headers):
        set_cell_text(insp_table.cell(0, i), h, bold=True, bg_color=grey_color)

    # Row 1 (UTTM)
    set_cell_text(insp_table.cell(1, 0), "UTTM", bold=True, bg_color=RGBColor(235, 241, 222)) # Light green tint
    set_cell_text(insp_table.cell(1, 1), "100% of TML Location – Refer TML as per Attachment", align=PP_ALIGN.LEFT, bg_color=RGBColor(235, 241, 222))
    set_cell_text(insp_table.cell(1, 2), "Cooling Water Corrosion", align=PP_ALIGN.LEFT, bg_color=RGBColor(235, 241, 222))

    # Row 2 (Visual)
    set_cell_text(insp_table.cell(2, 0), "VISUAL\nINSPECTION", bold=True, bg_color=RGBColor(235, 241, 222))
    set_cell_text(insp_table.cell(2, 1), "External Visual Inspection (100% Coverage)\n\nWith\n\nfollow up by UT, RT, or pit gauge as required.", align=PP_ALIGN.LEFT, bg_color=RGBColor(235, 241, 222))
    set_cell_text(insp_table.cell(2, 2), "General Corrosion", align=PP_ALIGN.LEFT, bg_color=RGBColor(235, 241, 222))


    # --- 5. DRAWING PLACEHOLDER (Bottom Right) ---
    # Add a rectangle or image placeholder
    left = Inches(5.8)
    top = Inches(3.8) # Slightly overlapping with inspection table area visually, aligned to drawing
    width = Inches(5.39)
    height = Inches(3.5)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = "[EQUIPMENT DRAWING PLACEHOLDER]"
    shape.fill.background() # Transparent fill
    line = shape.line
    line.color.rgb = RGBColor(180, 180, 180) # Light grey border

    # --- 6. FOOTER ---
    footer_top = Inches(7.5)
    footer_h = Inches(0.6)
    footer_table = slide.shapes.add_table(2, 3, Inches(0.5), footer_top, Inches(10.69), footer_h).table
    
    footer_headers = ["Prepared by", "Review and Approved by", "Acknowledge and Acc. by"]
    for i, h in enumerate(footer_headers):
        set_cell_text(footer_table.cell(0, i), h, bold=True, font_size=8, align=PP_ALIGN.LEFT)
        set_cell_text(footer_table.cell(1, i), "Date", bold=True, font_size=8, align=PP_ALIGN.CENTER)

    # Save
    prs.save(output_file)
    print(f"Template saved to {output_file}")
    return output_file

# --- Example Usage ---
# You can pass real data here, or leave it empty to see the placeholders
sample_data = {
    "tag_no": "V-101",
    "description": "Amine Flash Drum",
    "plant_unit": "01",
    "components": [
        {"name": "Top Head", "code": "ASME VIII DIV 1", "mat": "CS"},
        {"name": "Shell", "code": "ASME VIII DIV 1", "mat": "CS"},
        {"name": "Bottom Head", "code": "ASME VIII DIV 1", "mat": "CS"}
    ]
}

# Run the function
create_inspection_plan_template("Inspection_Plan_V101.pptx", sample_data)