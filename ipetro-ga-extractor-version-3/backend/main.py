from dotenv import load_dotenv
import os
from pathlib import Path
import base64
import io
import json as js
import re
import time
import uuid
from datetime import datetime, timedelta, timezone
from typing import List, Optional

# Load environment variables from .env file
env_path = Path(__file__).parent / ".env"
load_dotenv(dotenv_path=env_path)

import openai
import pandas as pd
from fastapi import (
    FastAPI,
    File,
    UploadFile,
    HTTPException,
    Depends,
    status,
    Form,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from sqlalchemy.orm import Session
import uuid
from google import genai
from google.genai import errors as genai_errors
from google.genai import types
from groq import Groq
from passlib.context import CryptContext
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pydantic import BaseModel, EmailStr
import jwt

from database import get_db, init_db
from database import (
    User as DBUser,
    Extraction as DBExtraction,
    Equipment as DBEquipment,
)
from google_sheets import GoogleSheetsClient

# --- configuration ---

app = FastAPI(title="RBI Equipment Manager", version="1.0.0")


# initialize database tables
@app.on_event("startup")
async def startup_event():
    init_db()
    print("✅ Database Initialized! ")
    print(
        f"DEBUG: GEMINI_KEY: {GEMINI_API_KEY[:5]}..."
        if GEMINI_API_KEY
        else "DEBUG: GEMINI_KEY is None/Empty"
    )
    print(
        f"DEBUG: OPENAI_KEY: {OPENAI_API_KEY[:5]}..."
        if OPENAI_API_KEY
        else "DEBUG: OPENAI_KEY is None/Empty"
    )
    print(
        f"DEBUG: GROQ_KEY: {GROQ_API_KEY[:5]}..."
        if GROQ_API_KEY
        else "DEBUG: GROQ_KEY is None/Empty"
    )


app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# folder
UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "outputs"
TEMPLATES_FOLDER = "templates"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
os.makedirs(TEMPLATES_FOLDER, exist_ok=True)

SECRET_KEY = os.environ.get("SECRET_KEY", "secret-key")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 1440  # 24 hours

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
security = HTTPBearer()


# AI providers setup
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
GROQ_API_KEY = os.environ.get("GROQ_API_KEY")

# initialize AI clients
gemini_client = genai.Client(api_key=GEMINI_API_KEY) if GEMINI_API_KEY else None
openai_client = openai.OpenAI(api_key=OPENAI_API_KEY) if OPENAI_API_KEY else None
groq_client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None

AVAILABLE_MODELS = {
    "gemini": [
        "gemini-2.5-flash",
        "gemini-1.5-pro",
        "gemini-1.5-flash",
        "gemini-3-pro-preview",
        "gemini-2.5-flash-lite",
        "gemini-3-flash-preview",
    ],
    "openai": ["gpt-4o", "gpt-4o-mini", "gpt-4-turbo"],
    "groq": ["llama-3.2-90b-vision-instruct", "llama-3.2-11b-vision-instruct"],
}

# in memory database (will replace with PostgreSQL/MongoDB in production)

users_db = {}
equipment_db = {}
extractions_db = {}

# Google Sheets Setup
SPREADSHEET_ID = "1RyzHPFmO9axcVc-one-K0Iy9IL4hXm1ucXzRicYlN5k"
SERVICE_ACCOUNT_FILE = "service_account.json"
gs_client = None

if os.path.exists(SERVICE_ACCOUNT_FILE):
    try:
        gs_client = GoogleSheetsClient(SERVICE_ACCOUNT_FILE, SPREADSHEET_ID)
        print("✅ Google Sheets Client Initialized")
    except Exception as e:
        print(f"❌ Failed to initialize Google Sheets Client: {e}")
else:
    print(f"❌ Service account file not found: {SERVICE_ACCOUNT_FILE}")


# --- pydantic models ---


class UserRegister(BaseModel):
    email: EmailStr
    password: str
    full_name: str
    company: Optional[str] = None


class UserLogin(BaseModel):
    email: EmailStr
    password: str


class Token(BaseModel):
    access_token: str
    token_type: str
    user: dict


class Equipment(BaseModel):
    tag: str
    type: str
    fluid: str
    material_type: str
    material_spec: str
    material_grade: str
    design_pressure: float
    design_temperature: float
    operating_pressure: float
    operating_temp: float
    insulation: bool


class EquipmentCreate(BaseModel):
    tag: str
    type: str
    fluid: str
    material_type: str
    material_spec: str
    material_grade: str
    design_pressure: float
    design_temperature: float
    operating_pressure: float
    operating_temp: float
    insulation: bool


class ExtractionResponse(BaseModel):
    id: str
    filename: str
    status: str
    provider: str
    model: str
    excel_file: Optional[str] = None
    pptx_file: Optional[str] = None
    created_at: datetime


# === authentication functions ===
def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)


def get_password_hash(password):
    return pwd_context.hash(password)


def create_access_token(data: dict):
    to_encode = data.copy()
    expire = datetime.now(timezone.utc) + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt


def decode_token(token: str):
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        return payload
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token has expired")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")


async def get_current_user(
    credentials: HTTPAuthorizationCredentials = Depends(security),
    db: Session = Depends(get_db),
):
    """get current authenticated user"""
    token = credentials.credentials
    payload = decode_token(token)
    email = payload.get("sub")

    if email is None:
        raise HTTPException(status_code=401, detail="Invalid authentication")

    user = db.query(DBUser).filter(DBUser.email == email).first()
    if user is None:
        raise HTTPException(status_code=401, detail="User not found")

    return user


# === AI processing functions ===
def allowed_file(filename: str) -> bool:
    ALLOWED_EXTENSIONS = {"pdf", "png", "jpg", "jpeg"}
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def encode_image_to_base64(image_bytes):
    return base64.b64encode(image_bytes).decode("utf-8")


def clean_text(text):
    return re.sub(r"```json|```", "", text, flags=re.IGNORECASE).strip()


def extract_json_from_text(text):
    match = re.search(r"(\{.*\}|\[.*\])", text, re.DOTALL)
    if match:
        return match.group(1).strip()
    return None


def call_gemini(image_bytes, prompt, model, max_retries=3):
    if not gemini_client:
        raise ValueError("Gemini API key is not configured.")

    contents = [
        types.Part.from_bytes(data=image_bytes, mime_type="image/jpeg"),
        prompt,
    ]

    for attempt in range(1, max_retries + 1):
        try:
            response = gemini_client.models.generate_content(
                model=model,
                contents=contents,
            )
            return response.text
        except genai_errors.ServerError as e:
            if e.code == 503 and attempt < max_retries:
                print(f"{model} is overloaded, (attempt {attempt}/{max_retries})")
                time.sleep(5)
            else:
                raise
    raise RuntimeError(f"{model} overloaded too many times.")


def call_openai(image_bytes, prompt, model):
    if not openai_client:
        raise ValueError("OpenAI API key is not configured.")

    base64_image = encode_image_to_base64(image_bytes)

    response = openai_client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                    },
                ],
            }
        ],
        max_tokens=4096,
    )

    return response.choices[0].message.content


def call_groq(image_bytes, prompt, model):
    if not groq_client:
        raise ValueError("Groq API key is not configured.")

    base64_image = encode_image_to_base64(image_bytes)

    response = groq_client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                    },
                ],
            }
        ],
        max_tokens=4096,
    )

    return response.choices[0].message.content


def process_single_page(
    image_or_bytes,
    page_num,
    total_pages,
    provider="gemini",
    model=None,
    source_type="image",
):
    if isinstance(image_or_bytes, Image.Image):
        buffer = io.BytesIO()
        image_or_bytes = image_or_bytes.convert("RGB")
        image_or_bytes.save(buffer, format="JPEG", quality=85)
        image_bytes = buffer.getvalue()
    else:
        image_bytes = image_or_bytes

    prompt = f"""Analyze this {source_type} (page {page_num} of {total_pages}) and extract equipment data for RBI Inspection.
    
    Extract the following information:

    1. **General Equipment Info**:
        - Equipment Tag/ID (e.g. V-001) - from the title of the file.
        - Equipment Description - Normally at the bottom right of the drawing, locate 'Subject:', 'Drawing Title:', or 'Title:'.
        - PMT Number - normally started by the word PMT, then followed by a few digits.
    
    2. **Component Details** (Extract for EACH part: Top Head, Shell, Bottom Head if applicable):
        - Phase (Liquid/Vapor/Two-Phase)
        - Fluid Name
        - Design Code: Normalize all variations of 'ASME VIII Div 1)
        - Materials: Type, Spec, Grade
    
    3. **Conditions**:
        - Design Pressure (MPa) & Temperature (C)
        - Operating Pressure (MPa) & Temperature (C)
        - Insulation (Yes/No)
    
    4. **Damage Mechanisms**:
        - List any active damage mechanisms mentioned (e.g. Cooling Water Corrosion, General Corrosion) or potential ones based on fluid/material.

    CRITICAL INSTRUCTION FOR COMPONENT NAMING & MAPPING (Output "name" must match strictly):
    - If specific data (Materials, Phase, Fluid) is labeled "TUBE SIDE":
        - You MUST assign this data to the component named "Channel" on it (e.g,. Top Channel, Bottom Channel) and "Head". 
        - If "Tube Bundle" is also present, it likely shares "TUBE SIDE" data or has its own.
    - If data is labeled "SHELL SIDE":
        - You MUST assign this data to the component named "Shell".
    - EXTRACT ONLY data for the PRIMARY equipment shown in the drawing (the one identified in `equipment_tag`). Do not mix in data from other equipment tables if they appear in the background.
    - Use EXACT standard names: "Tube Bundle", "Channel", "Shell", "Top Head", "Bottom Head", "Top Channel", "Bottom Channel", "Head" .
    - Do NOT output "Tube Side" or "Shell Side" as the component name. Use the mapped component names above.

    Return ONLY valid JSON with this structure:
    {{
        "equipment_tag": "Tag",
        "plant_unit": "Unit",
        "description": "Description",
        "pmt_no": "PMT No",
        "equipment_type": "Type",
        "damage_mechanisms": "List of mechanisms",
        "conditions": {{
            "design_pressure": 0.0,
            "design_temp": 0.0,
            "operating_pressure": 0.0,
            "operating_temp": 0.0,
            "insulation": false
        }},
        "components": [
            {{
                "name": "Component Name (e.g. Top Head, Shell, Bottom Head, Tube Bundle, Channel, Top Channel, Bottom Channel, Head)",
                "phase": "Phase (Liquid/Gas)",
                "fluid": "Fluid",
                "design_code": "Normalize all variation of 'ASME BOILER & PRESSURE VESSEL SECTION VIII DIV.1' or similar to just 'ASME VII DIV 1'.",
                "material": {{ "type": "Type", "spec": "Spec", "grade": "Grade" }}
            }}
        ]
    }}

    If no equipment data found, return: {{"error": "No equipment data found"}}
    """

    if not model:
        model = AVAILABLE_MODELS[provider][0]

    try:
        if provider == "gemini":
            full_output = call_gemini(image_bytes, prompt, model)
        elif provider == "openai":
            full_output = call_openai(image_bytes, prompt, model)
        elif provider == "groq":
            full_output = call_groq(image_bytes, prompt, model)
        else:
            raise ValueError(f"Unknown provider: {provider}")

        return clean_text(full_output)

    except Exception as e:
        print(f"Error calling {provider}: {str(e)}")
        raise


def process_pdf_as_image(pdf_path, provider="gemini", model=None):
    pages = convert_from_path(pdf_path, dpi=200)
    total_pages = len(pages)

    # process the first page for equipment data
    json_output = process_single_page(
        pages[0], 1, total_pages, provider, model, source_type="PDF page"
    )

    try:
        json_str = extract_json_from_text(json_output)
        if json_str:
            return js.loads(json_str)
    except:
        pass

    return {"error": "No data found"}


def process_image_file(image_path, provider="gemini", model=None):
    image = Image.open(image_path)
    json_output = process_single_page(image, 1, 1, provider, model, source_type="image")

    try:
        json_str = extract_json_from_text(json_output)
        if json_str:
            return js.loads(json_str)
    except:
        pass

    return {"error": "No data found"}


def save_equipment_to_excel(
    equipment_data: dict, output_file: str, template_path: str = None
):
    from openpyxl import load_workbook
    from openpyxl.styles import Border, Side, Alignment
    from openpyxl.cell.cell import MergedCell

    if output_file and os.path.exists(output_file):
        # Load existing Master File
        wb = load_workbook(output_file)
        ws = wb.active
    elif template_path and os.path.exists(template_path):
        wb = load_workbook(template_path)
        ws = wb.active
    else:
        # Fallback: Try to generate template using the existing script
        try:
            from create_excel_template import create_excel_template

            generated_template_path = create_excel_template()
            if generated_template_path and os.path.exists(generated_template_path):
                wb = load_workbook(generated_template_path)
                ws = wb.active
            else:
                raise FileNotFoundError("Created template not found")
        except Exception as e:
            # Absolute fallback if generation fails
            print(f"Warning: Could not load or create template: {e}")
            from openpyxl import Workbook

            wb = Workbook()
            ws = wb.active
            ws.append(["No Template Found - Data Appended Below"])
            # Create basic headers matching expected columns
            headers = [
                "NO.",
                "EQUIPMENT NO.",
                "PMT NO.",
                "EQUIPMENT DESCRIPTION",
                "PARTS",
                "PHASE",
                "FLUID",
                "MATERIAL TYPE",
                "MATERIAL SPEC",
                "MATERIAL GRADE",
                "INSULATION",
                "DESIGN TEMP",
                "DESIGN PRESSURE",
                "OPERATING TEMP",
                "OPERATING PRESSURE",
            ]
            ws.append(headers)

    # Locate start row (assuming row 6 based on template)
    # Correct logic for appending:
    # We iterate from row 6 down.
    # A row is occupied if Col 5 (Part Name) has content (it's never merged)
    # OR Col 1 has content (Number).
    start_row = 6
    while True:
        cell_obj_a = ws.cell(row=start_row, column=1)
        cell_a = str(cell_obj_a.value or "")
        cell_e = str(ws.cell(row=start_row, column=5).value or "")

        # 1. Check for Placeholder -> Overwrite it
        if "{{NO}}" in cell_a:
            break

        # 2. Check for Content
        # If Col 5 has text, it's a component row -> Occupied.
        # If Col 1 has text (Number), it's the start of a block -> Occupied.
        # If Col 1 is a MergedCell, it's part of a previous block -> Occupied.
        is_merged = isinstance(cell_obj_a, MergedCell)
        has_content = is_merged or cell_e.strip() or cell_a.strip()

        if has_content:
            start_row += 1  # Check next row
        else:
            # Empty -> Start here?
            # Double check next row too? No, let's assume if this row is clean, we can write.
            break

    # Define styles to match template
    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )
    center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)

    # Extract Data Helper
    def get_val(data, path, default=""):
        """Safely get nested dict value."""
        try:
            for key in path:
                data = data[key]
            return data if data is not None else default
        except (KeyError, TypeError):
            return default

    # Correct Mapping based on 3-row structure (Top Head, Shell, Bottom Head)
    # Row 1: Top Head + General Info
    # Row 2: Shell
    # Row 3: Bottom Head

    # Common Equipment Data
    eq_tag = equipment_data.get("equipment_tag", "")
    pmt = equipment_data.get("pmt_no", "")
    desc = equipment_data.get("description", "")
    insul = "Yes" if get_val(equipment_data, ["conditions", "insulation"]) else "No"

    des_temp = get_val(equipment_data, ["conditions", "design_temp"])
    des_press = get_val(equipment_data, ["conditions", "design_pressure"])
    op_temp = get_val(equipment_data, ["conditions", "operating_temp"])
    op_press = get_val(equipment_data, ["conditions", "operating_pressure"])

    # Components Data
    raw_comps = equipment_data.get("components", [])

    # Normalize to list
    component_list = []
    if isinstance(raw_comps, dict):
        for key in ["top_head", "shell", "bottom_head"]:
            if key in raw_comps:
                c = raw_comps[key]
                if "name" not in c:
                    c["name"] = key.replace("_", " ").title()
                component_list.append(c)
        for k, v in raw_comps.items():
            if k not in ["top_head", "shell", "bottom_head"]:
                if "name" not in v:
                    v["name"] = k
                component_list.append(v)
    elif isinstance(raw_comps, list):
        component_list = raw_comps

    # Use at least 3 rows to match standard template, but expand if more components
    num_rows = max(3, len(component_list))

    # Clean up existing merges in target range (to prevent 'MergedCell is read-only' error)
    # The template might have empty rows pre-merged.
    target_min_row = start_row
    target_max_row = start_row + num_rows - 1

    # Identify merges to remove
    merges_to_remove = []
    for merged_range in ws.merged_cells.ranges:
        # Check if overlap (simple check: if min_row or max_row is within target)
        # Or better: strictly if any part touches our working rows
        if not (
            merged_range.max_row < target_min_row
            or merged_range.min_row > target_max_row
        ):
            merges_to_remove.append(merged_range)

    for m in merges_to_remove:
        ws.unmerge_cells(str(m))

    # 1. Main Columns (Merged)
    # NO.
    ws.cell(row=start_row, column=1, value=start_row - 5)
    ws.merge_cells(
        start_row=start_row,
        start_column=1,
        end_row=start_row + num_rows - 1,
        end_column=1,
    )

    # EQUIPMENT NO.
    ws.cell(row=start_row, column=2, value=eq_tag)
    ws.merge_cells(
        start_row=start_row,
        start_column=2,
        end_row=start_row + num_rows - 1,
        end_column=2,
    )

    # PMT NO.
    ws.cell(row=start_row, column=3, value=pmt)
    ws.merge_cells(
        start_row=start_row,
        start_column=3,
        end_row=start_row + num_rows - 1,
        end_column=3,
    )

    # DESCRIPTION
    ws.cell(row=start_row, column=4, value=desc)
    ws.merge_cells(
        start_row=start_row,
        start_column=4,
        end_row=start_row + num_rows - 1,
        end_column=4,
    )
    ws.cell(row=start_row, column=4).alignment = center_align

    # 2. Component Rows
    for i in range(num_rows):
        current_row = start_row + i
        comp = component_list[i] if i < len(component_list) else {}

        # Part Name
        ws.cell(row=current_row, column=5, value=comp.get("name", ""))

        # Details
        ws.cell(row=current_row, column=6, value=comp.get("phase", ""))
        ws.cell(row=current_row, column=7, value=comp.get("fluid", ""))

        mat = comp.get("material", {})
        if isinstance(mat, dict):
            ws.cell(row=current_row, column=8, value=mat.get("type", ""))
            ws.cell(row=current_row, column=9, value=mat.get("spec", ""))
            ws.cell(row=current_row, column=10, value=mat.get("grade", ""))
        else:
            ws.cell(row=current_row, column=8, value=str(mat) if mat else "")

        # Conditions (apply to all rows if present)
        # Only if we have a component (or should we fill blanks? Template usually fills blanks)
        # Let's fill blanks if it's within the first 3 rows (standard), otherwise checking if comp exists.
        # But for 'insul'/'des_temp' which are equipment-wide, we might validly display them on every row.
        # However, standard practice is often to repeat or merge. Here we repeat as per original code.
        ws.cell(row=current_row, column=11, value=insul)
        ws.cell(row=current_row, column=12, value=des_temp)
        ws.cell(row=current_row, column=13, value=des_press)
        ws.cell(row=current_row, column=14, value=op_temp)
        ws.cell(row=current_row, column=15, value=op_press)

    # Apply Styling
    for r in range(start_row, start_row + num_rows):
        for c in range(1, 16):
            cell = ws.cell(row=r, column=c)
            cell.border = thin_border
            cell.alignment = center_align

    wb.save(output_file)
    return output_file


def parse_filename_info(filename: str) -> dict:
    """
    Extracts Equipment Tag and PMT Number from filename.
    Expected patterns:
      Tag: Letter-DigitDigitDigit (e.g. V-001)
      PMT: PMT followed by digits (e.g. PMT 10102)
    """
    info = {}

    # Regex for Tag: [A-Za-z]-[0-9]{3}
    tag_match = re.search(r"([A-Za-z]-[0-9]{3})", filename)
    if tag_match:
        info["equipment_tag"] = tag_match.group(1)

    # Regex for PMT: PMT\s*[0-9]+
    pmt_match = re.search(r"(PMT\s*[0-9]+)", filename, re.IGNORECASE)
    if pmt_match:
        # Standardize PMT format
        info["pmt_no"] = pmt_match.group(1).upper()

    return info


# save to inspection plan template
def save_equipment_to_template_pptx(
    equipment_data: dict, output_file: str, image_path: str = None
):

    data = equipment_data if equipment_data else {}

    # 1. Check if template exists, if not create it
    template_path = os.path.join(TEMPLATES_FOLDER, "inspection_plan_template.pptx")
    if not os.path.exists(template_path):
        print("Template not found, creating new precise template...")
        from create_inspection_plan_template import create_precise_template

        create_precise_template(template_path)

    # 2. Load Presentation
    prs = Presentation(template_path)
    slide = prs.slides[0]  # Single slide template

    # --- Helper to safely replace text in a run, preserving style if possible ---
    # --- Helper to safely replace text in a run, preserving style if possible ---
    def replace_text_in_text_frame(text_frame, replacements):
        for paragraph in text_frame.paragraphs:
            # 1. Fast path: Check individual runs first (preserves mixed formatting best)
            for run in paragraph.runs:
                for key, val in replacements.items():
                    if key in run.text:
                        run.text = run.text.replace(key, str(val))

            # 2. Slow path: Check full paragraph text for split tags (e.g. {{T + AG}})
            # If we find a key in the full text that wasn't fixed by the run loop,
            # we must consolidate runs.
            full_text = paragraph.text
            needs_replacement = False
            for key in replacements:
                if key in full_text:
                    needs_replacement = True
                    break

            if needs_replacement:
                # Perform replacement on the full string
                new_text = full_text
                for key, val in replacements.items():
                    if key in new_text:
                        new_text = new_text.replace(key, str(val))

                # capture style from the first run (heuristic)
                font_name = None
                font_size = None
                font_bold = None
                font_italic = None
                font_color_rgb = None

                if paragraph.runs:
                    r0 = paragraph.runs[0]
                    font_name = r0.font.name
                    font_size = r0.font.size
                    font_bold = r0.font.bold
                    font_italic = r0.font.italic
                    try:
                        font_color_rgb = r0.font.color.rgb
                    except:
                        pass  # Color might not be RGB (e.g. theme color)

                # Apply new text (this clears all runs and creates a single new one)
                paragraph.text = new_text

                # Re-apply style to the new single run
                if paragraph.runs:
                    new_run = paragraph.runs[0]
                    if font_name:
                        new_run.font.name = font_name
                    if font_size:
                        new_run.font.size = font_size
                    if font_bold is not None:
                        new_run.font.bold = font_bold
                    if font_italic is not None:
                        new_run.font.italic = font_italic
                    if font_color_rgb:
                        try:
                            new_run.font.color.rgb = font_color_rgb
                        except:
                            pass

    def replace_text_in_shape(shape, replacements):
        if not shape.has_text_frame:
            return
        replace_text_in_text_frame(shape.text_frame, replacements)

    def replace_text_in_table(table, replacements):
        for row in table.rows:
            for cell in row.cells:
                # Cells always have a text_frame property, no need to check has_text_frame
                replace_text_in_text_frame(cell.text_frame, replacements)

    # --- Prepare Data Map ---
    # Map extract keys to template placeholders (Found via inspection: {{PLANT_UNIT}}, {{FLUID1}}, etc.)
    replacements = {
        "{{PLANT_UNIT}}": data.get("plant_unit", ""),
        "{{DESC}}": data.get("description", ""),
        "{{TAG}}": data.get("equipment_tag", ""),
        "{{PMT}}": data.get("pmt_no", ""),
        # Common
        "{{RISK_RATING}}": "LOW",  # Default
    }

    # TODO: Connect to Inspection page.
    # TODO: Date 1 -> Prepared by:
    # TODO: Date 2 -> Review and Approved by:
    # TODO: Date 3 -> Acknowledge and Acc. by:
    # Date placeholders
    current_date = datetime.now().strftime("%d-%b-%y")
    replacements["{{DATE1}}"] = current_date
    replacements["{{DATE2}}"] = current_date
    replacements["{{DATE3}}"] = current_date

    # Component specific placeholders
    comps = data.get("components", {})

    # Helper to flatten component data
    def get_comp_val(comp_key, field, subfield=None):
        if not comps:
            return ""
        # Handle if comps is a list (fallback)
        comp = {}
        if isinstance(comps, list):
            name_map = {
                "top_head": "Top Head",
                "shell": "Shell",
                "bottom_head": "Bottom Head",
            }
            target = name_map.get(comp_key)
            for c in comps:
                if c.get("name") == target:
                    comp = c
                    break
        else:
            comp = comps.get(comp_key, {})

        if subfield:
            return comp.get(field, {}).get(subfield, "") or ""
        return comp.get(field, "") or ""

    # Use components exactly as they appear in the data (Master File order)
    component_list = data.get("components", [])

    # Ensure it is a list (handle legacy dict if necessary, though Master File returns list)
    if isinstance(component_list, dict):
        # Fallback only if dict: convert values to list
        component_list = list(component_list.values())

    # Iterate for up to 6 slots (Template support)
    print(
        f"DEBUG: Generating PPTX with {len(component_list)} components: {[c.get('name') for c in component_list]}"
    )
    for i in range(1, 7):
        idx = i - 1
        comp = component_list[idx] if idx < len(component_list) else {}

        # Determine name: use 'name' field, or fallback logic if empty but comp exists
        comp_name = comp.get("name", "")

        replacements[f"{{{{COMP{i}}}}}"] = comp_name
        replacements[f"{{{{FLUID{i}}}}}"] = comp.get("fluid", "")
        replacements[f"{{{{DC{i}}}}}"] = comp.get("design_code", "") or (
            "ASME VIII DIV 1" if comp_name else ""
        )

        mat = comp.get("material", {})
        # Handle if material is string or dict
        if isinstance(mat, str):
            # If AI returned string, just put it in TYPE and leave others
            replacements[f"{{{{TYPE{i}}}}}"] = mat
            replacements[f"{{{{SPEC{i}}}}}"] = ""
            replacements[f"{{{{GR{i}}}}}"] = ""
        else:
            replacements[f"{{{{TYPE{i}}}}}"] = mat.get("type", "")
            replacements[f"{{{{SPEC{i}}}}}"] = mat.get("spec", "")
            replacements[f"{{{{GR{i}}}}}"] = mat.get("grade", "")

        # Common conditions applied to all (unless we extract per-component later)
        # Only populate if component exists to keep template clean?
        # Actually template likely has lines for 1, 2, 3. If 2 is empty, we should probably clear INS2 etc.
        if comp_name:
            replacements[f"{{{{INS{i}}}}}"] = (
                "YES" if data.get("conditions", {}).get("insulation") else "NO"
            )
            replacements[f"{{{{TEMP{i}}}}}"] = str(
                data.get("conditions", {}).get("operating_temp", "")
            )
            replacements[f"{{{{PRESS{i}}}}}"] = str(
                data.get("conditions", {}).get("operating_pressure", "")
            )
        else:
            replacements[f"{{{{INS{i}}}}}"] = ""
            replacements[f"{{{{TEMP{i}}}}}"] = ""
            replacements[f"{{{{PRESS{i}}}}}"] = ""

    # --- Image Replacement Logic ---
    # Look for {{IMAGE_PLACEHOLDER}} in tables or textboxes
    def replace_image_placeholder(shape):
        found = False
        if shape.has_text_frame:
            if "{{IMAGE_PLACEHOLDER}}" in shape.text:
                found = True
                # Clear text
                shape.text_frame.text = ""
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if "{{IMAGE_PLACEHOLDER}}" in cell.text:
                        found = True
                        cell.text_frame.text = ""  # Clear

        if found:
            # Logic to insert image at shape/cell position?
            # Since it's hard to get absolute position of a cell,
            # let's assume the placeholder is its own shape or we just place image over it.
            # Better approach: If we find it, we get the shape's position (if it wasn't a table cell).
            pass

    # Simplified Image Logic:
    # Just find the shape with the placeholder, get its dimensions, delete it, and add the image there.
    # Note: `python-pptx` doesn't give easy cell coordinates.
    # But usually the drawing placeholder is a standalone shape (Rectangle).

    drawing_path = None
    # Potentially pass source image path in equipment_data if available?
    # For now, we don't have the source image path here easily unless we passed it.
    # But let's check if we can pass it.
    # In `upload_drawing` (line 978), we call `save_equipment_to_template_pptx`.
    # We should update the call signature to accept `image_path`.

    # --- Image Replacement Logic ---
    # Simplified Image Logic:
    # Just find the shape with the placeholder, get its dimensions, delete it, and add the image there.

    if image_path and os.path.exists(image_path):
        for slide in prs.slides:
            # Check shapes
            target_shape = None
            for shape in slide.shapes:
                if shape.has_text_frame and "{{IMAGE_PLACEHOLDER}}" in shape.text:
                    target_shape = shape
                    break
                # Also check tables if needed, but usually image placeholder is a text box

            if target_shape:
                # Capture geometry
                left = target_shape.left
                top = target_shape.top
                width = target_shape.width
                height = target_shape.height

                # Remove the placeholder text (or we could delete the shape, but text frame clear is safer if it has borders)
                # Actually, if we want to replace it with an image, we should probably add the image on top
                target_shape.text_frame.text = ""

                # Convert PDF to image if needed (already handled in extraction, but here we might have the original PDF)
                # If image_path is PDF, we need to convert it first.
                # Helper to get first page image
                final_image_path = image_path
                temp_image = None

                if image_path.lower().endswith(".pdf"):
                    try:
                        from pdf2image import convert_from_path

                        images = convert_from_path(image_path, last_page=1)
                        if images:
                            temp_image = os.path.join(OUTPUT_FOLDER, "temp_drawing.jpg")
                            images[0].save(temp_image, "JPEG")
                            final_image_path = temp_image
                    except Exception as e:
                        print(f"Failed to convert PDF for PPTX: {e}")
                        final_image_path = None

                if final_image_path:
                    try:
                        pic = slide.shapes.add_picture(
                            final_image_path, left, top, width=width, height=height
                        )
                        # Move to back? Or leave on top.
                        print(f"Inserted image from {final_image_path}")
                    except Exception as e:
                        print(f"Failed to insert image: {e}")

                if temp_image and os.path.exists(temp_image):
                    os.remove(temp_image)

    # --- Execute Replacement ---
    # Iterate all shapes (textboxes)
    for shape in slide.shapes:
        if shape.has_text_frame:
            replace_text_in_shape(shape, replacements)
        if shape.has_table:
            replace_text_in_table(shape.table, replacements)

    # Save
    prs.save(output_file)
    print(f"Template populated and saved: {output_file}")
    return output_file


# def save_equipment_to_pptx(
#     equipment_data: dict, output_file: str, template_path: str = None
# ):  # TODO: will customise to match the template

#     if template_path and os.path.exists(template_path):
#         prs = Presentation(template_path)
#     else:
#         prs = Presentation()
#         prs.slide_width = Inches(10)
#         prs.slide_height = Inches(7.5)

#     # add title slide
#     title_slide = prs.slides.add_slide(prs.slide_layouts[0])
#     title_slide.shapes.title.text = (
#         f"Inspection Plan: {equipment_data.get('equipment_tag', 'Unknown')}"
#     )

#     # add data slide
#     blank_slide = prs.slides.add_slide(prs.slide_layouts[6])

#     # add equipment details
#     left = Inches(1)
#     top = Inches(1.5)
#     width = Inches(8)
#     height = Inches(5)

#     txBox = blank_slide.shapes.add_textbox(left, top, width, height)
#     tf = txBox.text_frame
#     tf.text = f"Equipment Tag: {equipment_data.get('equipment_tag', '')}\n"
#     tf.text += f"Type: {equipment_data.get('equipment_type', '')}\n"
#     tf.text += f"Design Pressure: {equipment_data.get('design_pressure', '')}barg\n"
#     tf.text += f"Design Temperature: {equipment_data.get('design_temp', '')}ºC\n"
#     tf.text += f"Material: {equipment_data.get('shell_material', {}).get('spec', '')}\n"

#     prs.save(output_file)
#     return output_file

# API Routes - Authentication


@app.post("/api/auth/register", response_model=Token)
async def register(user: UserRegister, db: Session = Depends(get_db)):
    # --- register new user ---

    # check if user exists
    existing_user = db.query(DBUser).filter(DBUser.email == user.email).first()
    if existing_user:
        raise HTTPException(status_code=400, detail="Email already registered.")

    # create new user
    hashed_password = get_password_hash(user.password)
    user_id = str(uuid.uuid4())

    db_user = DBUser(
        id=user_id,
        email=user.email,
        hashed_password=hashed_password,
        full_name=user.full_name,
        company=user.company,
        created_at=datetime.now(timezone.utc),
    )

    db.add(db_user)
    db.commit()
    db.refresh(db_user)

    access_token = create_access_token(data={"sub": user.email})

    return {
        "access_token": access_token,
        "token_type": "bearer",
        "user": {
            "id": db_user.id,
            "email": db_user.email,
            "full_name": db_user.full_name,
            "company": db_user.company,
            "created_at": (
                db_user.created_at.isoformat() if db_user.created_at else None
            ),
        },
    }


@app.post("/api/auth/login", response_model=Token)
async def login(user: UserLogin, db: Session = Depends(get_db)):
    # --- login user ---
    db_user = db.query(DBUser).filter(DBUser.email == user.email).first()

    if not db_user or not verify_password(user.password, db_user.hashed_password):
        raise HTTPException(status_code=401, detail="Invalid email or password")

    access_token = create_access_token(data={"sub": user.email})

    # user_data = db_user.copy()
    # user_data.pop("hashed_password")

    return {
        "access_token": access_token,
        "token_type": "bearer",
        "user": {
            "id": db_user.id,
            "email": db_user.email,
            "full_name": db_user.full_name,
            "company": db_user.company,
            "created_at": db_user.created_at,
        },
    }


@app.get("/api/auth/me")
async def get_me(current_user: dict = Depends(get_current_user)):
    """get current user information"""
    return current_user


#  --- API Routes --- Models


@app.get("/api/models")
async def get_available_models():
    """get available models"""
    return {
        "models": AVAILABLE_MODELS,
        "configured": {
            "gemini": bool(GEMINI_API_KEY),
            "openai": bool(OPENAI_API_KEY),
            "groq": bool(GROQ_API_KEY),
        },
    }


# --- API Routes - Upload & Extraction ---
@app.post("/api/upload")
async def upload_drawing(
    file: UploadFile = File(...),
    provider: str = Form("gemini"),
    model: str = Form(None),
    current_user: DBUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """upload and extract equipment data from drawing"""
    if not allowed_file(file.filename):
        raise HTTPException(status_code=400, detail="Invalid file type")

    if provider == "gemini" and not GEMINI_API_KEY:
        raise HTTPException(status_code=400, detail="Gemini API key not configured")
    elif provider == "openai" and not OPENAI_API_KEY:
        raise HTTPException(status_code=400, detail="OpenAI API key not configured")
    elif provider == "groq" and not GROQ_API_KEY:
        raise HTTPException(status_code=400, detail="Groq API key not configured")

    # save file
    unique_id = str(uuid.uuid4())[:8]
    file_extension = Path(file.filename).suffix
    saved_filename = f"{unique_id}_{file.filename}"
    file_path = os.path.join(UPLOAD_FOLDER, saved_filename)

    with open(file_path, "wb") as buffer:
        content = await file.read()
        buffer.write(content)

    try:
        # process file
        if file.filename.lower().endswith(".pdf"):
            equipment_data = process_pdf_as_image(file_path, provider, model)
        else:
            equipment_data = process_image_file(file_path, provider, model)

        if "error" in equipment_data:
            raise HTTPException(status_code=400, detail=equipment_data["error"])

        # 2b. Parse info from filename and merge
        filename_info = parse_filename_info(file.filename)
        if filename_info:
            print(f"Extracted info from filename: {filename_info}")
            # Merge: filename info takes precedence or fills gaps
            # Here we override because user trusts filename more
            equipment_data.update(filename_info)

        # generate excel and pptx
        base_name = f"{unique_id}_{Path(file.filename).stem}"
        MASTER_FILENAME = "MasterFile_IPETRO_PLANT.xlsx"

        # Excel goes to Master File (persistent)
        excel_path = os.path.join(OUTPUT_FOLDER, MASTER_FILENAME)
        # PPTX stays unique per upload
        pptx_path = os.path.join(OUTPUT_FOLDER, f"{base_name}.pptx")

        template_excel = os.path.join(TEMPLATES_FOLDER, "equipment_template.xlsx")
        template_pptx = os.path.join(TEMPLATES_FOLDER, f"inspection_plan_template.pptx")

        if gs_client:
            try:
                gs_client.update_or_append_equipment_data(equipment_data)
                print("✅ Data synced with Google Sheet (Merge/Append)")
            except Exception as e:
                print(f"❌ Failed to sync with Google Sheet: {e}")
        else:
            # Fallback to local excel if GS fails or not configured?
            # For now, let's keep local excel as backup or just skip.
            # User wants to REPLACE it, but keeping local backup might be good.
            # Let's keep local generation for "download" purposes if they still want it.
            # But the requirement says "data push directly into masterfile... rather than download excel".
            # So the primary "MasterFile" is now the Google Sheet.
            # We can still save a local copy for the "excel_file" download link,
            # OR we should update the "excel_file" link to point to the Google Sheet URL?
            pass

        save_equipment_to_excel(
            equipment_data,
            excel_path,
            template_excel if os.path.exists(template_excel) else None,
        )
        save_equipment_to_template_pptx(
            equipment_data,
            pptx_path,
            image_path=file_path,
        )

        # save to database
        extraction_id = str(uuid.uuid4())
        db_extraction = DBExtraction(
            id=extraction_id,
            user_id=current_user.id,
            filename=file.filename,
            status="complete",
            provider=provider,
            model=model or AVAILABLE_MODELS[provider][0],
            equipment_data=equipment_data,
            excel_file=MASTER_FILENAME,
            pptx_file=f"{base_name}.pptx",
            completed_at=datetime.now(timezone.utc),
        )

        db.add(db_extraction)
        db.commit()
        db.refresh(db_extraction)

        # # clean up uploaded file
        # os.remove(file_path)

        # Upload PPTX to Drive if Client exists
        pptx_google_link = None
        if gs_client:
            pptx_google_link = gs_client.upload_pptx(pptx_path, f"{base_name}.pptx")
            if pptx_google_link:
                print(f"✅ PPTX uploaded to Drive: {pptx_google_link}")

        return {
            "success": True,
            "extraction_id": extraction_id,
            "equipment_data": equipment_data,
            "excel_file": MASTER_FILENAME,
            "pptx_file": f"{base_name}.pptx",
            "pptx_google_link": pptx_google_link,
        }
    except Exception as e:
        db.rollback()
        # Save failed extraction to database
        error_msg = str(e)
        print(f"Error processing upload: {error_msg}")

        extraction_id = str(uuid.uuid4())
        try:
            # Default model if provider lookup fails
            fallback_model = "unknown"
            if provider in AVAILABLE_MODELS and AVAILABLE_MODELS[provider]:
                fallback_model = AVAILABLE_MODELS[provider][0]

            db_extraction = DBExtraction(
                id=extraction_id,
                user_id=current_user.id,
                filename=file.filename,
                status="failed",
                provider=provider,
                model=model or fallback_model,
                error_message=error_msg,
                created_at=datetime.now(timezone.utc),
                completed_at=datetime.now(timezone.utc),
            )
            db.add(db_extraction)
            db.commit()
        except Exception as db_e:
            print(f"Failed to save error status to DB: {db_e}")

        # if os.path.exists(file_path):
        #     os.remove(file_path)
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/extractions")
async def list_extractions(
    current_user: dict = Depends(get_current_user), db: Session = Depends(get_db)
):
    """list all extractions for current user"""
    extractions = (
        db.query(DBExtraction)
        .filter(DBExtraction.user_id == current_user.id)
        .order_by(DBExtraction.created_at.desc())
        .all()
    )

    # Force UTC timezone if naive to ensure correct frontend display
    for ext in extractions:
        if ext.created_at and ext.created_at.tzinfo is None:
            ext.created_at = ext.created_at.replace(tzinfo=timezone.utc)

    return extractions


@app.get("/api/extractions/{extraction_id}")
async def get_extraction(
    extraction_id: str,
    current_user: DBUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """get specific extraction for current user"""
    extraction = (
        db.query(DBExtraction)
        .filter(
            DBExtraction.id == extraction_id, DBExtraction.user_id == current_user.id
        )
        .first()
    )

    if not extraction:
        raise HTTPException(status_code=404, detail="Extraction not found")

    if extraction.created_at and extraction.created_at.tzinfo is None:
        extraction.created_at = extraction.created_at.replace(tzinfo=timezone.utc)

    return extraction


# --- API routes - equipment ---


@app.post("/api/equipment")
async def create_equipment(
    equipment: EquipmentCreate,
    current_user: DBUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """register new equipment manually"""
    equipment_id = str(uuid.uuid4())

    db_equipment = DBEquipment(
        id=equipment_id,
        user_id=current_user.id,
        **equipment.dict(),
        created_at=datetime.now(timezone.utc),
    )

    db.add(db_equipment)
    db.commit()
    db.refresh(db_equipment)

    return db_equipment


@app.get("/api/equipment")
async def list_equipment(
    current_user: DBUser = Depends(get_current_user), db: Session = Depends(get_db)
):
    """list all equipment for current user"""
    equipment = (
        db.query(DBEquipment)
        .filter(DBEquipment.user_id == current_user.id)
        .order_by(DBEquipment.created_at.desc())
        .all()
    )

    return equipment


@app.get("/api/equipment/{equipment_id}")
async def get_equipment(
    equipment_id: str,
    current_user: DBUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """get specific equipment for current user"""
    equipment = (
        db.query(DBEquipment)
        .filter(DBEquipment.id == equipment_id, DBEquipment.user_id == current_user.id)
        .first()
    )

    if not equipment:
        raise HTTPException(status_code=404, detail="Equipment not found")

    return equipment


@app.get("/api/download/{file_type}/{filename}")
async def download_file(
    file_type: str, filename: str, current_user: dict = Depends(get_current_user)
):
    """download pptx or excel file"""
    file_path = os.path.join(OUTPUT_FOLDER, filename)

    if not os.path.exists(file_path):
        raise HTTPException(status_code=400, detail="File not found")

    if file_type == "excel":
        media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    elif file_type == "pptx":
        media_type = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    else:
        raise HTTPException(status_code=400, detail="Invalid file type")

    return FileResponse(file_path, media_type=media_type, filename=filename)


# --- health check ---


@app.get("/")
async def root():
    return {"message": "RBI Equipment API", "version": "1.0,0", "status": "running"}


@app.get("/health")
async def health_check():
    return {
        "status": "healthy",
        "ai_providers": {
            "gemini": bool(GEMINI_API_KEY),
            "openai": bool(OPENAI_API_KEY),
            "groq": bool(GROQ_API_KEY),
        },
    }


# Run with : uvicorn main:app --reload --port 8000


# --- Master File Features ---


def read_equipment_from_master_file(tag: str, master_file_path: str):
    """
    Reads equipment data from the Master Excel file for a specific tag.
    Handles merged cells (Equipment info merged across multiple rows).
    Returns a dict structure compatible with save_equipment_to_template_pptx.
    """
    from openpyxl import load_workbook

    if not os.path.exists(master_file_path):
        return None

    wb = load_workbook(master_file_path, data_only=True)
    ws = wb.active

    # Find rows for the tag
    # Iterate from row 6
    max_row = ws.max_row

    start_row = 6

    # First pass: Identify the block of rows for this tag
    current_tag = None
    target_block_start = -1
    target_block_end = -1

    for r in range(start_row, max_row + 2):  # +2 to ensure we process the last block
        cell_val = ws.cell(row=r, column=2).value  # Equipment No

        # If cell has a value, it starts a new block (or is a single row)
        if cell_val:
            # If we were tracking a tag, this new value ends the previous block
            if current_tag == tag:
                target_block_end = r - 1
                break  # Found our block finish

            # Start tracking new tag
            current_tag = str(cell_val).strip()
            if current_tag == tag:
                target_block_start = r

        # Check if we are inside the target block but reached end of file/data
        # (e.g. next ID appears, or EOF)
        if current_tag == tag and r == max_row:
            target_block_end = max_row
            break

    if target_block_start == -1:
        return None  # Tag not found

    if target_block_end == -1:
        target_block_end = target_block_start  # Single row case if logic fell through

    # Extract Data from the block
    # General Info (from first row of block)
    row_1 = target_block_start
    equipment_data = {
        "equipment_tag": tag,
        "pmt_no": ws.cell(row=row_1, column=3).value,
        "description": ws.cell(row=row_1, column=4).value,
        "conditions": {
            # Accessing from first row (assuming merged/repeated)
            "insulation": ws.cell(row=row_1, column=11).value == "Yes",  # Col 11
            "design_temp": ws.cell(row=row_1, column=12).value,
            "design_pressure": ws.cell(row=row_1, column=13).value,
            "operating_temp": ws.cell(row=row_1, column=14).value,
            "operating_pressure": ws.cell(row=row_1, column=15).value,
        },
        "components": [],
    }

    # Extract Components (Iterate through the block)
    for r in range(target_block_start, target_block_end + 1):
        part_name = ws.cell(row=r, column=5).value
        # If part name is empty, it might be a truly empty row inside a block, skip
        if not part_name:
            continue

        comp = {
            "name": part_name,
            "phase": ws.cell(row=r, column=6).value,
            "fluid": ws.cell(row=r, column=7).value,
            "material": {
                "type": ws.cell(row=r, column=8).value,
                "spec": ws.cell(row=r, column=9).value,
                "grade": ws.cell(row=r, column=10).value,
            },
        }
        equipment_data["components"].append(comp)

    return equipment_data


class GeneratePPTXRequest(BaseModel):
    tag: str


@app.post("/api/generate-pptx")
async def generate_pptx_endpoint(
    request: GeneratePPTXRequest, current_user: DBUser = Depends(get_current_user)
):
    if not gs_client:
        raise HTTPException(
            status_code=500, detail="Google Sheets Client not initialized."
        )

    try:
        data = gs_client.get_equipment_data(request.tag)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Google Sheets Error: {str(e)}")

    if not data:
        raise HTTPException(
            status_code=404,
            detail=f"Equipment Tag '{request.tag}' not found in Master Google Sheet.",
        )

    # Generate PPTX
    base_name = f"Generated_{request.tag}"
    pptx_path = os.path.join(OUTPUT_FOLDER, f"{base_name}.pptx")

    try:
        # We need an image path? The original logic passed `uploaded_file_path`.
        # If generating from master file, we might not have the image handy easily unless we stored its path in DB or MasterFile.
        # For now, we'll pass None or a default placeholder.
        # The user's request implies we are regenerating from DATA.
        save_equipment_to_template_pptx(data, pptx_path, image_path=None)
    except Exception as e:
        print(f"Error generating PPTX: {e}")
        raise HTTPException(
            status_code=500, detail=f"Failed to generate PPTX: {str(e)}"
        )

    return {
        "success": True,
        "pptx_file": f"{base_name}.pptx",
        "message": f"PPTX generated for {request.tag}",
    }
