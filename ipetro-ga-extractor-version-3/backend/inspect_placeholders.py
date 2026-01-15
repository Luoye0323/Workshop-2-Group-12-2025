from pptx import Presentation
import os


def inspect_pptx(file_path):
    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        return

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error loading PPTX: {e}")
        return

    print(f"Inspecting: {file_path}")
    print("-" * 30)

    found_placeholders = set()

    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}:")

        # Check shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text
                    if "{{" in text and "}}" in text:
                        # Simple extraction of {{KEY}}
                        import re

                        matches = re.findall(r"\{\{.*?\}\}", text)
                        for m in matches:
                            found_placeholders.add(m)
                            print(f"  Found placeholder: {m} (in shape)")

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                text = paragraph.text
                                if "{{" in text and "}}" in text:
                                    import re

                                    matches = re.findall(r"\{\{.*?\}\}", text)
                                    for m in matches:
                                        found_placeholders.add(m)
                                        print(f"  Found placeholder: {m} (in table)")

    print("-" * 30)
    print(f"Unique placeholders found: {sorted(list(found_placeholders))}")


if __name__ == "__main__":
    target_file = "/Users/muhammadputraazam/Desktop/uni/S5/workshop/Workshop-2-Group-12-2025/ipetro-ga-extractor-version-3/backend/templates/inspection_plan_template.pptx"
    inspect_pptx(target_file)
