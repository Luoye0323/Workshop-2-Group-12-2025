import base64
import io
import json as js
import os
import re
import time
import uuid
from pathlib import Path

import openai
import pandas as pd
import xlsxwriter
from flask import Flask, jsonify, render_template, request, send_file, session
from flask_cors import CORS
# ai provider imports
from google import genai
from google.genai import errors as genai_errors
from google.genai import types
from groq import Groq
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from werkzeug.utils import secure_filename

app = Flask(__name__)
CORS(app)
app.config["SECRET_KEY"] = "secret key"  # TODO secret keyyy
app.config["UPLOAD_FOLDER"] = "uploads"
app.config["OUTPUT_FOLDER"] = "outputs"
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  # 50mb
app.config["ALLOWED_EXTENSIONS"] = {"pdf", "png", "jpg", "jpeg"}

# create folders if they do not exist
os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)
os.makedirs(app.config["OUTPUT_FOLDER"], exist_ok=True)

# --- api provider setup ---
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "")
GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "")

# initialize clients
gemini_client = genai.Client(api_key=GEMINI_API_KEY) if GEMINI_API_KEY else None
openai_client = openai.OpenAI(api_key=OPENAI_API_KEY) if OPENAI_API_KEY else None
groq_client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None

# available models per provider
AVAILABLE_MODELS = {
    "gemini": ["gemini-2.5-flash", "gemini-1.5-pro", "gemini-1.5-flash"],
    "openai": ["gpt-4o", "gpt-4o-mini", "gpt-4-turbo"],
    "groq": ["llama-3.2-90b-vision-preview", "llama-3.2-11b-vision-preview"],
}


def allowed_file(filename):
    return (
        "." in filename
        and filename.rsplit(".", 1)[1].lower() in app.config["ALLOWED_EXTENSIONS"]
    )


def encode_image_to_base64(image_bytes):
    """convert image bytes to base64 string"""
    return base64.b64decode(image_bytes).decode("utf-8")


def clean_text(text):
    return re.sub(r"```json|```", "", text, flags=re.IGNORECASE).strip()


def extract_json_from_text(text):
    match = re.search(r"(\{.*\}|\[.*\])", text, re.DOTALL)
    if match:
        return match.group(1).strip()
    return None


def call_gemini(image_bytes, prompt, model="gemini-2.5-flash", max_retries=3):
    if not gemini_client:
        raise ValueError("Gemini API key is not configured.")

    contents = [
        types.Part.from_bytes(data=image_bytes, mime_type="image/jpeg"),
        prompt,
    ]

    for attempt in range(1, max_retries + 1):
        try:
            responses = gemini_client.models.generate_content_stream(
                model=model, contents=contents
            )
            full_output = ""
            for response in responses:
                if response.text:
                    full_output += response.text
            return full_output
        except genai_errors.ServerError as e:
            if e.code == 503 and attempt < max_retries:
                print(
                    f"⚠️ {model} overloaded (attempt {attempt}/{max_retries}), retrying..."
                )
                time.sleep(5)
            else:
                raise

    raise RuntimeError(f"{model} overloaded too many times. Terminating...")


def call_openai(image_bytes, prompt, model="gpt-4o"):
    if not openai_client:
        raise ValueError("OpenAI API key is not configured.")

    base64_image = encode_image_to_base64(image_bytes)

    response = openai_client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                    },
                ],
            }
        ],
        max_tokens=4096,
    )

    return response.choices[0].message.content


def call_groq(image_bytes, prompt, model="llama-3.2-90b-vision-preview"):
    if not groq_client:
        raise ValueError("Groq API is not configured.")

    base64_image = encode_image_to_base64(image_bytes)

    response = groq_client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                    },
                ],
            }
        ],
        max_tokens=4096,
    )

    return response.choices[0].message.content


def process_single_page(
    image_or_bytes,
    page_num,
    total_pages,
    provider="gemini",
    model=None,
    source_type="image",
):
    """process a single page with selected ai provider"""

    # convert image to bytes - if needed
    if isinstance(image_or_bytes, Image.Image):
        buffer = io.BytesIO()
        image_or_bytes = image_or_bytes.convert("RGB")
        image_or_bytes.save(buffer, format="JPEG", quality=85)
        image_bytes = buffer.getvalue()
    else:
        image_bytes = image_or_bytes

        # prompts - TODO: edit for better extraction... maybe after this.
        prompt = f"""Analyze this {source_type} (page {page_num}/{total_pages}) and extract ALL tables found.

            INSTRUCTIONS:
            1. Detect ALL tables in this page (there may be multiple tables)
            2. For EACH table found, extract complete data
            3. Extract EVERY single cell - do not skip or omit any data
            4. Preserve the EXACT table structure (all rows and columns)
            5. Include ALL headers exactly as they appear
            6. Include ALL rows - even if cells are empty, represent them as empty strings ""
            7. Maintain column alignment - each row must have the same number of columns as headers
            8. Include numbers, text, symbols, and special characters exactly as shown
            9. Do not merge, summarize, or truncate any data

            If NO tables are found, return: {{"tables": []}}

            If tables ARE found, return JSON in this EXACT format:
            {{
            "tables": [
                {{
                "table_number": 1,
                "headers": ["col1", "col2", "col3", ...],
                "rows": [["val1", "val2", "val3", ...], ["val1", "val2", "val3", ...], ...]
                }}
            ]
            }}

            Return ONLY valid JSON - no markdown, no explanations, no additional text."""

        # set to default model (gemini) if not specified by the user
        if not model:
            model = AVAILABLE_MODELS[provider][0]  # <- gemini

        # --- call appropiate provider (gemini, openai, groq)
        try:
            if provider == "gemini":
                full_output = call_gemini(image_bytes, prompt, model)
            elif provider == "openai":
                full_output = call_openai(image_bytes, prompt, model)
            elif provider == "groq":
                full_output = call_groq(image_bytes, prompt, model)
            else:
                raise ValueError(f"Unknown provider: {provider}")

            return clean_text(full_output)

        except Exception as e:
            print(f"Error calling {provider}: {str(e)}")
            raise


def process_pdf_as_image(pdf_path, provider="gemini", model=None):
    """process pdf file with selected ai provider"""
    pages = convert_from_path(pdf_path, dpi=200)
    total_pages = len(pages)
    all_tables = []

    for page_num, page_image in enumerate(pages, 1):
        json_output = process_single_page(
            page_image, page_num, total_pages, provider, model, source_type="PDF page"
        )

        try:
            json_str = extract_json_from_text(json_output)
            if json_str:
                data = js.load(json_str)
                if "tables" in data and data["tables"]:
                    for table in data["tables"]:
                        table["page"] = page_num
                        all_tables.append(table)
        except Exception as e:
            print(f"Error parsing page {page_num}: {e}")

        if page_num < total_pages:
            time.sleep(2)

    return {"tables": all_tables}


def process_image_file(image_path, provider="gemini", model=None):
    """process image file with selected ai provider"""
    image = Image.open(image_path)
    json_output = process_single_page(image, 1, 1, provider, model, source_type="image")

    try:
        json_str = extract_json_from_text(json_output)
        if json_str:
            return js.loads(json_str)
    except:
        pass

    return {"tables": []}


def save_tables_to_excel(tables_data, output_file):
    """save table to excel file"""
    tables = tables_data.get("tables", [])

    if not tables:
        return None

    with pd.ExcelWriter(output_file, engine="xlsxwriter") as writer:
        workbook = writer.book
        header_format = workbook.add_format(
            {"bold": True, "bg_color": "#DCE6F1", "border": 1}
        )

        for idx, table in enumerate(tables, 1):
            page = table.get("page", "1")
            table_num = table.get("table_number", idx)

            if "header" not in table or "rows" not in table:
                continue

            headers = table["headers"]
            rows = table["rows"]

            # fixing the length of the rows
            fixed_rows = []
            for row in rows:
                row = (row + [""] * len(headers))[: len(headers)]
                fixed_rows.append(row)

            df = pd.DataFrame(fixed_rows, columns=headers)
            sheet_name = f"Page{page}_Table{table_num}"[:31]

            df.to_excel(writer, index=False, sheet_name=sheet_name)
            worksheet = writer.sheets[sheet_name]

            for col_num, value in enumerate(df.columns.values):
                worksheet.write(0, col_num, value, header_format)

    return output_file


def save_tables_to_pptx(tables_data, output_file, provider):
    """save tables to powerpoint file
    TODO: need to confugurate the template to sync with given templates by ipetro"""

    tables = tables_data.get("tables", [])

    if not tables:
        return None

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholder[1]
    title.text = "Extracted Tables"
    subtitle.text = f"Total Tables: {len(tables)} by {provider}"

    for idx, table_data in enumerate(tables, 1):
        page = table_data.get("page", "?")
        table_num = table_data.get("table_number", idx)

        if "headers" not in table_data or "rows" not in table_data:
            continue

        headers = table_data["headers"]
        rows = table_data["rows"]

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"Page {page} - Table {table_num}"
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.bold = True

        num_cols = len(headers)
        num_rows = len(rows) + 1

        table_width = Inches(9)
        table_height = min(Inches(5.5), Inches(0.4 * num_rows))

        for col_idx, header in enumerate(headers):
            cell = tables.cell(0, col_idx)
            cell.text = str(header)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)

        for row_idx, row in enumerate(rows, 1):
            for col_idx, cell_value in enumerate(row):
                if col_idx < num_cols:
                    cell = tables.cell(row_idx, col_idx)
                    cell.text = str(cell_value) if cell_value else ""
                    cell.text_frame.paragraphs[0].font.size = Pt(9)

    prs.save(output_file)
    return output_file


# --- Flask ---


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/api/models", methods=["GET"])
def get_available_models():
    """return available models and their configuration status"""
    return jsonify(
        {
            "models": AVAILABLE_MODELS,
            "configured": {
                "gemini": bool(GEMINI_API_KEY),
                "openai": bool(OPENAI_API_KEY),
                "groq": bool(GROQ_API_KEY),
            },
        }
    )


@app.route("/api/upload", methods=["POST"])
def upload_file():
    try:
        if "files[]" not in request.files:
            return jsonify({"error": "No files provided"}), 400

        files = request.files.getlist("files[]")
        provider = request.form.get("provider", "gemini")
        model = request.form.get("model", None)

        if not files or files[0].filename == "":
            return jsonify({"error": "No files selected"}), 400

        # validating provider if it configured
        if provider == "gemini" and not GEMINI_API_KEY:
            return jsonify({"error": "Gemini API key not configured"}), 400
        elif provider == "openai" and not OPENAI_API_KEY:
            return jsonify({"error": "OpenAI API key not configured"})
        elif provider == "groq" and not GROQ_API_KEY:
            return jsonify({"error": "Groq API key not configured"})

        results = []

        for file in files:
            if file and allowed_file(file.filename):
                unique_id = str(uuid.uuid4())[:8]
                filename = secure_filename(file.filename)
                file_path = os.path.join(
                    app.config["UPLOAD_FOLDER"], f"{unique_id}_{filename}"
                )
                file.save(file_path)

            try:
                if filename.lower().endswith(".pdf"):
                    tables_data = process_pdf_as_image(file_path, provider, model)
                else:
                    tables_data = process_image_file(file_path, provider, model)

                num_tables = len(tables_data.get("tables", []))

                if num_tables > 0:
                    base_name = f"{unique_id}_{Path(filename).stem}"
                    excel_path = os.path.join(
                        app.config["OUTPUT_FOLDER"], f"{base_name}.xlsx"
                    )
                    pptx_path = os.path.join(
                        app.config["OUTPUT_FOLDER"], f"{base_name}.pptx"
                    )

                    save_tables_to_excel(tables_data, excel_path)
                    save_tables_to_pptx(tables_data, pptx_path)

                    results.append(
                        {
                            "filename": filename,
                            "status": "success",
                            "tables_found": num_tables,
                            "excel_file": f"{base_name}.xlsx",
                            "pptx_file": f"{base_name}.pptx",
                            "provider": provider,
                            "model": model,
                        }
                    )
                else:
                    results.append(
                        {"filename": filename, "status": "no_tables", "tables_found": 0}
                    )

                os.remove(file_path)
            except Exception as e:
                print(f"Error processing {filename}: {str(e)}")
                import traceback

                traceback.print_exc()
                results.append(
                    {"filename": filename, "status": "error", "error": str(e)}
                )
        else:
            results.append(
                {
                    "filename": file.filename if file else "unknown",
                    "status": "invalid_type",
                }
            )
    finally:
        return jsonify({"error": f"Server error : {str(e)}"}), 500


@app.route("/api/download/<file_type>/<filename>")
def download_file(file_type, filename):
    file_path = os.path.join(app.config["OUTPUT_FOLDER"], filename)

    if os.path.exists(file_path):
        if file_path == "excel":
            mimetype = (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )
        elif file_path == "pptx":
            mimetype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        else:
            return "Invalid file type", 400

        return send_file(
            file_path, mimetype=mimetype, as_attachment=True, download_name=filename
        )

    return "File not found", 404

@app.route("/api/test/json", methods=['GET'])
def return_home():
    return jsonify({
        'message': "Hello testing"
    })


if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5001)
