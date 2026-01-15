# app.py - Flask Web Application for Table Extraction

import io
import json as js
import os
import re
import time
import uuid
from pathlib import Path

from flask_cors import CORS
from flask_sqlalchemy import SQLAlchemy
from werkzeug.security import generate_password_hash, check_password_hash
import jwt
from datetime import datetime, timedelta
from functools import wraps

import pandas as pd
import xlsxwriter
from flask import Flask, jsonify, render_template, request, send_file, session
from google import genai
from google.genai import errors as genai_errors
from google.genai import types
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from werkzeug.utils import secure_filename


app = Flask(__name__)

CORS(app, resources={
    r"/api/*": {
        "origins": ["http://localhost:3000"],
        "methods": ["GET", "POST", "PUT", "DELETE"],
        "allow_headers": ["Content-Type", "Authorization"]
    }
})

app.config['SECRET_KEY'] = 'your-secret-key-here-change-this-in-production'
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'outputs'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'png', 'jpg', 'jpeg'}

# Initialize database - FIXED: Added missing slash
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///tables_app.db'  # Fixed: was 'sqlite://tables_app.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

db = SQLAlchemy(app)

# Create folders if they don't exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

# === Gemini Setup ===
os.environ["GEMINI_API_KEY"] = "AIzaSyCUN3KvVvx57LtX27awuK_jdCmVPtTtP4w"
client = genai.Client(api_key=os.environ["GEMINI_API_KEY"])


# ===== DATABASE MODELS =====

class User(db.Model):
    """User model for authentication"""
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password_hash = db.Column(db.String(200), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

    # Relationship to files
    files = db.relationship('ProcessedFile', backref='user', lazy=True)

    def set_password(self, password):
        self.password_hash = generate_password_hash(password)

    def check_password(self, password):
        return check_password_hash(self.password_hash, password)


class ProcessedFile(db.Model):
    """Store information about processed files"""
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    original_filename = db.Column(db.String(200), nullable=False)
    unique_id = db.Column(db.String(50), nullable=False)
    tables_found = db.Column(db.Integer, default=0)
    excel_file = db.Column(db.String(200))
    pptx_file = db.Column(db.String(200))
    processed_at = db.Column(db.DateTime, default=datetime.utcnow)
    status = db.Column(db.String(50), default='success')


# Create database tables
with app.app_context():
    db.create_all()


# ===== AUTHENTICATION DECORATOR =====

def token_required(f):
    """Decorator to protect routes - requires valid JWT token"""
    @wraps(f)
    def decorated(*args, **kwargs):
        token = None

        # Get token from Authorization header
        if 'Authorization' in request.headers:
            auth_header = request.headers['Authorization']
            try:
                token = auth_header.split(" ")[1]  # FIXED: was split(" "),[1] with comma
            except IndexError:
                return jsonify({'error': 'Invalid token format'}), 401

        if not token:
            return jsonify({'error': 'Token is missing'}), 401  # FIXED: was 'token' not 'error'
        
        try:
            # Decode token
            data = jwt.decode(token, app.config['SECRET_KEY'], algorithms=["HS256"])
            current_user = User.query.get(data['user_id'])
            if not current_user:
                return jsonify({'error': 'User not found'}), 401
        except jwt.ExpiredSignatureError:
            return jsonify({'error': 'Token has expired'}), 401
        except jwt.InvalidTokenError:
            return jsonify({'error': 'Invalid token'}), 401
        
        return f(current_user, *args, **kwargs)
    
    return decorated


# ===== AUTHENTICATION ROUTES =====

@app.route('/api/auth/register', methods=['POST'])
def register():
    """Register new user - connects to REGISTER BUTTON"""
    try:
        data = request.get_json()

        # Validation
        if not data.get('username') or not data.get('email') or not data.get('password'):
            return jsonify({'error': 'Missing required fields'}), 400

        # Check if user already exists
        if User.query.filter_by(username=data['username']).first():
            return jsonify({'error': 'Username already exists'}), 400
        
        if User.query.filter_by(email=data['email']).first():  # FIXED: was filer_by (typo)
            return jsonify({'error': 'Email already exists'}), 400

        # Create new user
        new_user = User(
            username=data['username'],
            email=data['email']
        )
        new_user.set_password(data['password'])

        db.session.add(new_user)
        db.session.commit()

        return jsonify({
            'message': 'User registered successfully',
            'user': {
                'id': new_user.id,
                'username': new_user.username,
                'email': new_user.email
            }
        }), 201
    
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


@app.route('/api/auth/login', methods=['POST'])
def login():
    """Login user - connects to LOGIN BUTTON"""
    try:
        data = request.get_json()

        if not data.get('username') or not data.get('password'):
            return jsonify({'error': 'Username and password required'}), 400
        
        # Find user
        user = User.query.filter_by(username=data['username']).first()

        if not user or not user.check_password(data['password']):
            return jsonify({'error': 'Invalid username or password'}), 401  # FIXED: was {'Invalid...'} with extra braces
        
        # Create JWT token - FIXED: Multiple errors here
        token = jwt.encode({
            'user_id': user.id,
            'exp': datetime.utcnow() + timedelta(hours=24)
        }, app.config['SECRET_KEY'], algorithm="HS256")  # FIXED: was ==, removed extra ), and ==

        return jsonify({
            'message': 'Login successful',  # FIXED: typo 'successfull'
            'token': token,
            'user': {
                'id': user.id,
                'username': user.username,
                'email': user.email
            }
        }), 200
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/auth/me', methods=['GET'])
@token_required
def get_current_user(current_user):
    """Get current logged-in user info"""
    return jsonify({
        'id': current_user.id,
        'username': current_user.username,
        'email': current_user.email,
        'created_at': current_user.created_at.isoformat()
    })


# ===== HELPER FUNCTIONS =====

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']


def safe_generate(client, model, contents, max_retries=3, delay=5):
    for attempt in range(1, max_retries + 1):
        try:
            return client.models.generate_content_stream(model=model, contents=contents)
        except genai_errors.ServerError as e:
            if e.code == 503:
                print(f"⚠️ {model} overloaded (attempt {attempt}/{max_retries}), retrying in {delay}s...")
                time.sleep(delay)
            else:
                raise
    raise RuntimeError(f"❌ {model} overloaded too many times.")


def clean_text(text):
    return re.sub(r"```json|```", "", text, flags=re.IGNORECASE).strip()


def extract_json_from_text(text):
    match = re.search(r'(\{.*\}|\[.*\])', text, re.DOTALL)
    if match:
        return match.group(1).strip()
    return None


def process_single_page(image_or_bytes, page_num, total_pages, source_type="image"):
    if isinstance(image_or_bytes, Image.Image):
        buffer = io.BytesIO()
        image_or_bytes = image_or_bytes.convert("RGB")
        image_or_bytes.save(buffer, format="JPEG", quality=85)
        image_bytes = buffer.getvalue()
    else:
        image_bytes = image_or_bytes
    
    prompt = f"""Analyze this {source_type} (page {page_num}/{total_pages}) and extract ALL tables found.

INSTRUCTIONS:
1. Detect ALL tables in this page (there may be multiple tables)
2. For EACH table found, extract complete data
3. Extract EVERY single cell - do not skip or omit any data
4. Preserve the EXACT table structure (all rows and columns)
5. Include ALL headers exactly as they appear
6. Include ALL rows - even if cells are empty, represent them as empty strings ""
7. Maintain column alignment - each row must have the same number of columns as headers
8. Include numbers, text, symbols, and special characters exactly as shown
9. Do not merge, summarize, or truncate any data

If NO tables are found, return: {{"tables": []}}

If tables ARE found, return JSON in this EXACT format:
{{
  "tables": [
    {{
      "table_number": 1,
      "headers": ["col1", "col2", "col3", ...],
      "rows": [["val1", "val2", "val3", ...], ["val1", "val2", "val3", ...], ...]
    }}
  ]
}}

Return ONLY valid JSON - no markdown, no explanations, no additional text."""
    
    contents = [
        types.Part.from_bytes(data=image_bytes, mime_type="image/jpeg"),
        prompt,
    ]

    try:
        responses = safe_generate(client, "gemini-2.5-flash", contents)
    except RuntimeError:
        responses = safe_generate(client, "gemini-1.5-pro", contents)
    
    full_output = ""
    for response in responses:
        if response.text:
            full_output += response.text

    return clean_text(full_output)


def process_pdf_as_images(pdf_path):
    pages = convert_from_path(pdf_path, dpi=200)
    total_pages = len(pages)
    all_tables = []
    
    for page_num, page_image in enumerate(pages, 1):
        json_output = process_single_page(page_image, page_num, total_pages, source_type="PDF page")
        
        try:
            json_str = extract_json_from_text(json_output)
            if json_str:
                data = js.loads(json_str)
                if "tables" in data and data["tables"]:
                    for table in data["tables"]:
                        table["page"] = page_num
                        all_tables.append(table)
        except Exception as e:
            print(f"Error parsing page {page_num}: {e}")
        
        if page_num < total_pages:
            time.sleep(2)
    
    return {"tables": all_tables}


def process_image_file(image_path):
    image = Image.open(image_path)
    json_output = process_single_page(image, 1, 1, source_type="image")
    
    try:
        json_str = extract_json_from_text(json_output)
        if json_str:
            return js.loads(json_str)
    except:
        pass
    
    return {"tables": []}


def save_tables_to_excel(tables_data, output_file):
    """Save tables to Excel file"""
    tables = tables_data.get("tables", [])
    
    if not tables:
        return None
    
    with pd.ExcelWriter(output_file, engine="xlsxwriter") as writer:
        workbook = writer.book
        header_format = workbook.add_format({
            "bold": True,
            "bg_color": "#DCE6F1",
            "border": 1
        })
        
        for idx, table in enumerate(tables, 1):
            page = table.get("page", "1")
            table_num = table.get("table_number", idx)
            
            if "headers" not in table or "rows" not in table:
                continue
            
            headers = table["headers"]
            rows = table["rows"]
            
            # Fix row lengths
            fixed_rows = []
            for row in rows:
                row = (row + [""] * len(headers))[:len(headers)]
                fixed_rows.append(row)
            
            df = pd.DataFrame(fixed_rows, columns=headers)
            sheet_name = f"Page{page}_Table{table_num}"[:31]
            
            df.to_excel(writer, index=False, sheet_name=sheet_name)
            worksheet = writer.sheets[sheet_name]
            
            for col_num, value in enumerate(df.columns.values):
                worksheet.write(0, col_num, value, header_format)
    
    return output_file


def save_tables_to_pptx(tables_data, output_file):
    """Save tables to PowerPoint file"""
    tables = tables_data.get("tables", [])
    
    if not tables:
        return None
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Extracted Tables"
    subtitle.text = f"Total Tables: {len(tables)}"
    
    # Add each table as a slide
    for idx, table_data in enumerate(tables, 1):
        page = table_data.get("page", "?")
        table_num = table_data.get("table_number", idx)
        
        if "headers" not in table_data or "rows" not in table_data:
            continue
        
        headers = table_data["headers"]
        rows = table_data["rows"]
        
        # Create blank slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = f"Page {page} - Table {table_num}"
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.bold = True
        
        # Calculate table size
        num_cols = len(headers)
        num_rows = len(rows) + 1
        
        table_width = Inches(9)
        table_height = min(Inches(5.5), Inches(0.4 * num_rows))
        
        # Add table
        table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1), 
                                      table_width, table_height).table
        
        # Fill headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        # Fill data rows
        for row_idx, row in enumerate(rows, 1):
            for col_idx, cell_value in enumerate(row):
                if col_idx < num_cols:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_value) if cell_value else ""
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
    
    prs.save(output_file)
    return output_file


# ===== FILE PROCESSING ROUTES =====

@app.route('/api/upload', methods=['POST'])
@token_required
def upload_file(current_user):
    """Upload and process files - connects to UPLOAD BUTTON"""
    try:
        if 'files[]' not in request.files:
            return jsonify({'error': 'No files provided'}), 400
        
        files = request.files.getlist('files[]')
        
        if not files or files[0].filename == '':
            return jsonify({'error': 'No files selected'}), 400
        
        results = []
        
        for file in files:
            if file and allowed_file(file.filename):
                # Generate unique filename
                unique_id = str(uuid.uuid4())[:8]
                filename = secure_filename(file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_{filename}")
                file.save(file_path)
                
                try:
                    # Process file
                    if filename.lower().endswith('.pdf'):
                        tables_data = process_pdf_as_images(file_path)
                    else:
                        tables_data = process_image_file(file_path)
                    
                    num_tables = len(tables_data.get("tables", []))
                    
                    if num_tables > 0:
                        # Generate output files
                        base_name = f"{unique_id}_{Path(filename).stem}"
                        excel_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{base_name}.xlsx")
                        pptx_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{base_name}.pptx")
                        
                        save_tables_to_excel(tables_data, excel_path)
                        save_tables_to_pptx(tables_data, pptx_path)

                        # Save to database
                        processed_file = ProcessedFile(
                            user_id=current_user.id,
                            original_filename=filename,
                            unique_id=unique_id,
                            tables_found=num_tables,
                            excel_file=f"{base_name}.xlsx",
                            pptx_file=f"{base_name}.pptx",
                            status='success'
                        )
                        db.session.add(processed_file)
                        db.session.commit()
                        
                        results.append({
                            'id': processed_file.id,
                            'filename': filename,
                            'status': 'success',
                            'tables_found': num_tables,
                            'excel_file': f"{base_name}.xlsx",
                            'pptx_file': f"{base_name}.pptx"
                        })
                    else:
                        results.append({
                            'filename': filename,
                            'status': 'no_tables',
                            'tables_found': 0
                        })
                    
                    # Clean up uploaded file
                    os.remove(file_path)
                    
                except Exception as e:
                    print(f"❌ Error processing {filename}: {str(e)}")
                    import traceback
                    traceback.print_exc()
                    results.append({
                        'filename': filename,
                        'status': 'error',
                        'error': str(e)
                    })
            else:
                results.append({
                    'filename': file.filename if file else 'unknown',
                    'status': 'invalid_type'
                })
        
        return jsonify({'results': results})
    
    except Exception as e:
        print(f"❌ Upload route error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'Server error: {str(e)}'}), 500


@app.route('/api/files', methods=['GET'])
@token_required
def get_user_files(current_user):
    """Get all files processed by current user - connects to FILE HISTORY"""
    files = ProcessedFile.query.filter_by(user_id=current_user.id).order_by(
        ProcessedFile.processed_at.desc()
    ).all()

    return jsonify({
        'files': [{
            'id': f.id,
            'filename': f.original_filename,
            'tables_found': f.tables_found,
            'excel_file': f.excel_file,
            'pptx_file': f.pptx_file,
            'processed_at': f.processed_at.isoformat(),
            'status': f.status
        } for f in files]
    })


@app.route('/api/download/<file_type>/<filename>')
@token_required
def download_file(current_user, file_type, filename):
    """Download generated files - connects to DOWNLOAD BUTTON"""
    # Verify file belongs to user
    file_record = ProcessedFile.query.filter_by(
        user_id=current_user.id
    ).filter(
        (ProcessedFile.excel_file == filename) | (ProcessedFile.pptx_file == filename)
    ).first()

    if not file_record:
        return jsonify({'error': 'File not found or access denied'}), 404

    file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    
    if os.path.exists(file_path):
        if file_type == 'excel':
            mimetype = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        elif file_type == 'pptx':
            mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        else:
            return jsonify({'error': 'Invalid file type'}), 400
        
        return send_file(file_path, mimetype=mimetype, as_attachment=True, download_name=filename)
    
    return jsonify({'error': 'File not found'}), 404


if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5001)